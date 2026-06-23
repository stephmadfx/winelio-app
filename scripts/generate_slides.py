import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    c_orange = RGBColor(255, 107, 53)     # #FF6B35 (Orange Winelio)
    c_amber = RGBColor(247, 147, 30)      # #F7931E (Ambre Winelio)
    c_charcoal = RGBColor(45, 52, 54)     # #2D3436 (Dark charcoal)
    c_slate = RGBColor(61, 74, 82)        # #3D4A52 (Slate grey)
    c_grey = RGBColor(99, 110, 114)       # #636E72 (Neutral grey)
    c_light = RGBColor(248, 249, 250)     # #F8F9FA (Off-white)
    c_white = RGBColor(255, 255, 255)
    
    # File paths for logos
    logo_dir = "/Users/steph/PROJETS/WINELIO/winelio/logo/png"
    logo_color = os.path.join(logo_dir, "winelio-logo-color.png")
    logo_on_dark = os.path.join(logo_dir, "winelio-logo-on-dark.png")
    logo_white = os.path.join(logo_dir, "winelio-logo-white.png")
    logo_tagline_white = os.path.join(logo_dir, "winelio-logo-tagline-white.png")
    
    # Helper function to set slide background color
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper function to create text frame with common defaults
    def create_text_box(slide, left, top, width, height):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        return tf

    blank_layout = prs.slide_layouts[6] # completely blank layout

    # ==========================================
    # SLIDE 1: Page de Garde (Dark)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_background(s1, c_charcoal)
    
    # Add Winelio Logo
    if os.path.exists(logo_on_dark):
        s1.shapes.add_picture(logo_on_dark, Inches(4.666), Inches(1.8), width=Inches(4.0))
    
    # Title
    tf_title = create_text_box(s1, Inches(1.0), Inches(3.6), Inches(11.333), Inches(1.2))
    p_title = tf_title.paragraphs[0]
    p_title.text = "La Plateforme de Recommandation Professionnelle"
    p_title.font.name = "Poppins"
    p_title.font.size = Pt(38)
    p_title.font.bold = True
    p_title.font.color.rgb = c_white
    p_title.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    tf_sub = create_text_box(s1, Inches(1.0), Inches(4.7), Inches(11.333), Inches(0.8))
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Réseau MLM & Recommandations d'Affaires de Confiance"
    p_sub.font.name = "Poppins"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = c_amber
    p_sub.alignment = PP_ALIGN.CENTER
    
    # Bottom orange/amber accent bar
    accent_bar1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.3), Inches(6.666), Inches(0.2))
    accent_bar1.fill.solid()
    accent_bar1.fill.fore_color.rgb = c_orange
    accent_bar1.line.fill.background()
    
    accent_bar2 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.666), Inches(7.3), Inches(6.667), Inches(0.2))
    accent_bar2.fill.solid()
    accent_bar2.fill.fore_color.rgb = c_amber
    accent_bar2.line.fill.background()

    # Presenter tag
    tf_pres = create_text_box(s1, Inches(1.0), Inches(6.4), Inches(11.333), Inches(0.5))
    p_pres = tf_pres.paragraphs[0]
    p_pres.text = "Présentation de l'Application — Juin 2026"
    p_pres.font.name = "Poppins"
    p_pres.font.size = Pt(12)
    p_pres.font.color.rgb = c_grey
    p_pres.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 2: Le Concept (Light)
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_background(s2, c_light)
    
    # Top-right Logo
    if os.path.exists(logo_color):
        s2.shapes.add_picture(logo_color, Inches(10.5), Inches(0.5), width=Inches(2.0))
        
    # Title
    tf_s2_t = create_text_box(s2, Inches(0.8), Inches(0.5), Inches(9.0), Inches(0.8))
    p_s2_t = tf_s2_t.paragraphs[0]
    p_s2_t.text = "Une nouvelle façon de connecter"
    p_s2_t.font.name = "Poppins"
    p_s2_t.font.size = Pt(28)
    p_s2_t.font.bold = True
    p_s2_t.font.color.rgb = c_charcoal
    
    # Orange line decoration under title
    dec_line = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.06))
    dec_line.fill.solid()
    dec_line.fill.fore_color.rgb = c_orange
    dec_line.line.fill.background()
    
    # Content Columns
    # Left Column (Concept)
    tf_col1 = create_text_box(s2, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
    p1 = tf_col1.paragraphs[0]
    p1.text = "La Recommandation Simplifiée"
    p1.font.name = "Poppins"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = c_orange
    p1.space_after = Pt(14)
    
    bullets1 = [
        "Mise en relation de confiance entre particuliers et professionnels certifiés.",
        "Suivi transparent des dossiers à travers un workflow structuré en 8 étapes.",
        "Une application moderne disponible sur mobile et desktop pour un suivi en temps réel."
    ]
    for b in bullets1:
        p = tf_col1.add_paragraph()
        p.text = "• " + b
        p.font.name = "Poppins"
        p.font.size = Pt(15)
        p.font.color.rgb = c_slate
        p.space_after = Pt(12)
        
    # Right Column (Sponsor/MLM)
    tf_col2 = create_text_box(s2, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5))
    p2 = tf_col2.paragraphs[0]
    p2.text = "La Puissance du MLM"
    p2.font.name = "Poppins"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = c_amber
    p2.space_after = Pt(14)
    
    bullets2 = [
        "Inscription par code de parrainage obligatoire pour garantir la qualité du réseau.",
        "Structure de récompense sur 5 niveaux d'affiliation directs et indirects.",
        "Rémunération équitable basée sur la réussite réelle des projets de recommandation."
    ]
    for b in bullets2:
        p = tf_col2.add_paragraph()
        p.text = "• " + b
        p.font.name = "Poppins"
        p.font.size = Pt(15)
        p.font.color.rgb = c_slate
        p.space_after = Pt(12)

    # Footer slide number / brand
    tf_foot2 = create_text_box(s2, Inches(0.8), Inches(6.8), Inches(11.733), Inches(0.3))
    p_foot2 = tf_foot2.paragraphs[0]
    p_foot2.text = "Winelio App  |  Recommandez. Connectez. Développez."
    p_foot2.font.name = "Poppins"
    p_foot2.font.size = Pt(10)
    p_foot2.font.color.rgb = c_grey

    # ==========================================
    # SLIDE 3: Transition / MLM (Dark)
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_background(s3, c_charcoal)
    
    # Left vertical orange accent bar
    v_bar = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.08), Inches(3.0))
    v_bar.fill.solid()
    v_bar.fill.fore_color.rgb = c_orange
    v_bar.line.fill.background()
    
    # Big title text
    tf_s3 = create_text_box(s3, Inches(1.2), Inches(2.2), Inches(10.0), Inches(3.0))
    p_s3_1 = tf_s3.paragraphs[0]
    p_s3_1.text = "PARTIE II"
    p_s3_1.font.name = "Poppins"
    p_s3_1.font.size = Pt(14)
    p_s3_1.font.bold = True
    p_s3_1.font.color.rgb = c_orange
    p_s3_1.space_after = Pt(12)
    
    p_s3_2 = tf_s3.add_paragraph()
    p_s3_2.text = "Le Modèle MLM &\nLe Plan de Commissionnement"
    p_s3_2.font.name = "Poppins"
    p_s3_2.font.size = Pt(36)
    p_s3_2.font.bold = True
    p_s3_2.font.color.rgb = c_white
    p_s3_2.space_after = Pt(16)
    
    p_s3_3 = tf_s3.add_paragraph()
    p_s3_3.text = "Découvrez comment l'application rémunère son réseau d'apporteurs d'affaires sur 5 niveaux."
    p_s3_3.font.name = "Poppins"
    p_s3_3.font.size = Pt(16)
    p_s3_3.font.color.rgb = c_grey

    # Small white logo in corner
    if os.path.exists(logo_white):
        s3.shapes.add_picture(logo_white, Inches(10.5), Inches(0.5), width=Inches(2.0))

    # ==========================================
    # SLIDE 4: Commission Distribution (Light with 3 Cards)
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_background(s4, c_light)
    
    # Top-right Logo
    if os.path.exists(logo_color):
        s4.shapes.add_picture(logo_color, Inches(10.5), Inches(0.5), width=Inches(2.0))
        
    # Title
    tf_s4_t = create_text_box(s4, Inches(0.8), Inches(0.5), Inches(9.0), Inches(0.8))
    p_s4_t = tf_s4_t.paragraphs[0]
    p_s4_t.text = "Distribution d'une Commission"
    p_s4_t.font.name = "Poppins"
    p_s4_t.font.size = Pt(28)
    p_s4_t.font.bold = True
    p_s4_t.font.color.rgb = c_charcoal
    
    dec_line4 = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.06))
    dec_line4.fill.solid()
    dec_line4.fill.fore_color.rgb = c_orange
    dec_line4.line.fill.background()

    # Let's draw 3 Columns/Cards
    card_width = Inches(3.64)
    card_height = Inches(4.3)
    card_y = Inches(1.8)
    gap = Inches(0.4)
    
    # Card 1: L'apporteur d'affaires (60%)
    c1_left = Inches(0.8)
    card1 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c1_left, card_y, card_width, card_height)
    card1.fill.solid()
    card1.fill.fore_color.rgb = c_white
    card1.line.color.rgb = c_orange
    card1.line.width = Pt(1.5)
    
    tf_card1 = create_text_box(s4, c1_left + Inches(0.3), card_y + Inches(0.3), card_width - Inches(0.6), card_height - Inches(0.6))
    pc1_1 = tf_card1.paragraphs[0]
    pc1_1.text = "L'APPORTEUR"
    pc1_1.font.name = "Poppins"
    pc1_1.font.size = Pt(14)
    pc1_1.font.bold = True
    pc1_1.font.color.rgb = c_grey
    pc1_1.space_after = Pt(10)
    
    pc1_2 = tf_card1.add_paragraph()
    pc1_2.text = "60%"
    pc1_2.font.name = "Poppins"
    pc1_2.font.size = Pt(44)
    pc1_2.font.bold = True
    pc1_2.font.color.rgb = c_orange
    pc1_2.space_after = Pt(12)
    
    bulletsc1 = [
        "Reçoit la part majoritaire de la commission d'apport d'affaires.",
        "Rôle actif dans la mise en relation et l'initialisation du contact.",
        "Gains débloqués dès la validation finale du devis par le client."
    ]
    for b in bulletsc1:
        p = tf_card1.add_paragraph()
        p.text = "• " + b
        p.font.name = "Poppins"
        p.font.size = Pt(12)
        p.font.color.rgb = c_slate
        p.space_after = Pt(8)

    # Card 2: Le réseau de Parrainage (15%)
    c2_left = c1_left + card_width + gap
    card2 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c2_left, card_y, card_width, card_height)
    card2.fill.solid()
    card2.fill.fore_color.rgb = c_white
    card2.line.color.rgb = c_amber
    card2.line.width = Pt(1.5)
    
    tf_card2 = create_text_box(s4, c2_left + Inches(0.3), card_y + Inches(0.3), card_width - Inches(0.6), card_height - Inches(0.6))
    pc2_1 = tf_card2.paragraphs[0]
    pc2_1.text = "LE RÉSEAU MLM"
    pc2_1.font.name = "Poppins"
    pc2_1.font.size = Pt(14)
    pc2_1.font.bold = True
    pc2_1.font.color.rgb = c_grey
    pc2_1.space_after = Pt(10)
    
    pc2_2 = tf_card2.add_paragraph()
    pc2_2.text = "15% total"
    pc2_2.font.name = "Poppins"
    pc2_2.font.size = Pt(44)
    pc2_2.font.bold = True
    pc2_2.font.color.rgb = c_amber
    pc2_2.space_after = Pt(12)
    
    bulletsc2 = [
        "Réparti sur 5 niveaux d'affiliation au-dessus du recommandeur.",
        "3% de commission par niveau (Niveau 1 à Niveau 5).",
        "Encourage la formation, le coaching et le recrutement qualitatif."
    ]
    for b in bulletsc2:
        p = tf_card2.add_paragraph()
        p.text = "• " + b
        p.font.name = "Poppins"
        p.font.size = Pt(12)
        p.font.color.rgb = c_slate
        p.space_after = Pt(8)

    # Card 3: Plateforme & Bonus (25%)
    c3_left = c2_left + card_width + gap
    card3 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c3_left, card_y, card_width, card_height)
    card3.fill.solid()
    card3.fill.fore_color.rgb = c_white
    card3.line.color.rgb = c_slate
    card3.line.width = Pt(1.5)
    
    tf_card3 = create_text_box(s4, c3_left + Inches(0.3), card_y + Inches(0.3), card_width - Inches(0.6), card_height - Inches(0.6))
    pc3_1 = tf_card3.paragraphs[0]
    pc3_1.text = "PLATEFORME & BONUS"
    pc3_1.font.name = "Poppins"
    pc3_1.font.size = Pt(14)
    pc3_1.font.bold = True
    pc3_1.font.color.rgb = c_grey
    pc3_1.space_after = Pt(10)
    
    pc3_2 = tf_card3.add_paragraph()
    pc3_2.text = "25% total"
    pc3_2.font.name = "Poppins"
    pc3_2.font.size = Pt(44)
    pc3_2.font.bold = True
    pc3_2.font.color.rgb = c_slate
    pc3_2.space_after = Pt(12)
    
    bulletsc3 = [
        "23% pour le fonctionnement de la plateforme Winelio.",
        "1% Bonus d'affiliation pour le parrain du professionnel.",
        "1% Cashback Professionnel crédité en monnaie virtuelle (Wins)."
    ]
    for b in bulletsc3:
        p = tf_card3.add_paragraph()
        p.text = "• " + b
        p.font.name = "Poppins"
        p.font.size = Pt(12)
        p.font.color.rgb = c_slate
        p.space_after = Pt(8)

    # ==========================================
    # SLIDE 5: Workflow 8 Étapes (Light)
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_background(s5, c_light)
    
    # Top-right Logo
    if os.path.exists(logo_color):
        s5.shapes.add_picture(logo_color, Inches(10.5), Inches(0.5), width=Inches(2.0))
        
    # Title
    tf_s5_t = create_text_box(s5, Inches(0.8), Inches(0.5), Inches(9.0), Inches(0.8))
    p_s5_t = tf_s5_t.paragraphs[0]
    p_s5_t.text = "Workflow de Recommandation"
    p_s5_t.font.name = "Poppins"
    p_s5_t.font.size = Pt(28)
    p_s5_t.font.bold = True
    p_s5_t.font.color.rgb = c_charcoal
    
    dec_line5 = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.06))
    dec_line5.fill.solid()
    dec_line5.fill.fore_color.rgb = c_orange
    dec_line5.line.fill.background()

    # Let's layout 8 steps in a clean grid (2 rows of 4 columns)
    step_w = Inches(2.7)
    step_h = Inches(2.0)
    grid_gap_x = Inches(0.3)
    grid_gap_y = Inches(0.3)
    grid_x = Inches(0.8)
    grid_y = Inches(1.8)
    
    steps = [
        ("1. Recommandation reçue", "Le professionnel reçoit la notification d'apport d'affaires."),
        ("2. Acceptée par le pro", "Le professionnel accepte la recommandation et prend connaissance du lead."),
        ("3. Contact établi", "Prise de contact initiée avec le client potentiel (prospect)."),
        ("4. Rendez-vous fixé", "Planification d'une rencontre ou d'un appel technique pour qualifier le besoin."),
        ("5. Devis soumis", "Étude technique effectuée et envoi officiel de la proposition tarifaire."),
        ("6. Devis validé", "Signature du devis et versement de l'acompte. Déclenchement des commissions !"),
        ("7. Paiement reçu", "Règlement des prestations par le client selon l'échéancier convenu."),
        ("8. Affaire terminée", "Clôture du dossier et notation de la recommandation par les parties.")
    ]
    
    for idx, (title, desc) in enumerate(steps):
        row = idx // 4
        col = idx % 4
        x = grid_x + col * (step_w + grid_gap_x)
        y = grid_y + row * (step_h + grid_gap_y)
        
        # Step box
        box = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, step_w, step_h)
        box.fill.solid()
        box.fill.fore_color.rgb = c_white
        # Highlight step 6 which is critical (triggering commissions)
        if idx == 5:
            box.line.color.rgb = c_orange
            box.line.width = Pt(2.0)
        else:
            box.line.color.rgb = RGBColor(220, 225, 230)
            box.line.width = Pt(1.0)
            
        tf_step = create_text_box(s5, x + Inches(0.15), y + Inches(0.15), step_w - Inches(0.3), step_h - Inches(0.3))
        
        p_st = tf_step.paragraphs[0]
        p_st.text = title
        p_st.font.name = "Poppins"
        p_st.font.size = Pt(13)
        p_st.font.bold = True
        if idx == 5:
            p_st.font.color.rgb = c_orange
        else:
            p_st.font.color.rgb = c_charcoal
        p_st.space_after = Pt(6)
        
        p_sd = tf_step.add_paragraph()
        p_sd.text = desc
        p_sd.font.name = "Poppins"
        p_sd.font.size = Pt(10.5)
        p_sd.font.color.rgb = c_grey

    # ==========================================
    # SLIDE 6: Conclusion (Dark)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_background(s6, c_charcoal)
    
    # Logo in center
    if os.path.exists(logo_on_dark):
        s6.shapes.add_picture(logo_on_dark, Inches(4.666), Inches(1.8), width=Inches(4.0))

    # Tagline
    tf_tag = create_text_box(s6, Inches(1.0), Inches(3.6), Inches(11.333), Inches(0.8))
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = "Recommandez. Connectez. "
    p_tag.font.name = "Poppins"
    p_tag.font.size = Pt(24)
    p_tag.font.color.rgb = c_white
    p_tag.alignment = PP_ALIGN.CENTER
    
    # Highlight "Développez." in orange
    run = p_tag.add_run()
    run.text = "Développez."
    run.font.name = "Poppins"
    run.font.bold = True
    run.font.color.rgb = c_orange
    
    # Text info
    tf_info = create_text_box(s6, Inches(1.0), Inches(4.7), Inches(11.333), Inches(1.5))
    p_i1 = tf_info.paragraphs[0]
    p_i1.text = "Merci pour votre attention"
    p_i1.font.name = "Poppins"
    p_i1.font.size = Pt(20)
    p_i1.font.bold = True
    p_i1.font.color.rgb = c_white
    p_i1.alignment = PP_ALIGN.CENTER
    p_i1.space_after = Pt(10)
    
    p_i2 = tf_info.add_paragraph()
    p_i2.text = "Démo en direct disponible sur : dev2.winelio.app\nContact et support : contact@aide-multimedia.fr"
    p_i2.font.name = "Poppins"
    p_i2.font.size = Pt(14)
    p_i2.font.color.rgb = c_grey
    p_i2.alignment = PP_ALIGN.CENTER

    # Bottom orange/amber accent bar
    accent_bar1_6 = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.3), Inches(6.666), Inches(0.2))
    accent_bar1_6.fill.solid()
    accent_bar1_6.fill.fore_color.rgb = c_orange
    accent_bar1_6.line.fill.background()
    
    accent_bar2_6 = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.666), Inches(7.3), Inches(6.667), Inches(0.2))
    accent_bar2_6.fill.solid()
    accent_bar2_6.fill.fore_color.rgb = c_amber
    accent_bar2_6.line.fill.background()

    # Save presentation
    output_path = "/Users/steph/PROJETS/WINELIO/winelio/Winelio_Presentation_Template.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
