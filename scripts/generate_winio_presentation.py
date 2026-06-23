import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    c_orange = RGBColor(255, 107, 53)     # #FF6B35 (Orange Winelio)
    c_amber = RGBColor(247, 147, 30)      # #F7931E (Ambre Winelio)
    c_charcoal = RGBColor(45, 52, 54)     # #2D3436 (Dark charcoal)
    c_slate = RGBColor(61, 74, 82)        # #3D4A52 (Slate grey)
    c_grey = RGBColor(99, 110, 114)       # #636E72 (Neutral grey)
    c_light = RGBColor(248, 249, 250)     # #F8F9FA (Off-white)
    c_white = RGBColor(255, 255, 255)
    
    # File paths for logos
    logo_dir = "/Users/steph/PROJETS/WINELIO/winelio/logo/png"
    logo_color = os.path.join(logo_dir, "winelio-logo-color.png")
    logo_on_dark = os.path.join(logo_dir, "winelio-logo-on-dark.png")
    logo_white = os.path.join(logo_dir, "winelio-logo-white.png")
    
    # Helper to set slide background color
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to create text frame with margins reset
    def create_text_box(slide, left, top, width, height):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        return tf

    # Helper to add standard header to light slides
    def add_slide_header(slide, title_text):
        # Top-right Logo
        if os.path.exists(logo_color):
            slide.shapes.add_picture(logo_color, Inches(10.5), Inches(0.5), width=Inches(2.0))
            
        # Title
        tf = create_text_box(slide, Inches(0.8), Inches(0.5), Inches(9.0), Inches(0.8))
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Poppins"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = c_charcoal
        
        # Orange line decoration under title
        dec_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.06))
        dec_line.fill.solid()
        dec_line.fill.fore_color.rgb = c_orange
        dec_line.line.fill.background()

    # Helper to add standard footer to light slides
    def add_slide_footer(slide, current_slide, total_slides=9):
        tf = create_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.733), Inches(0.3))
        p = tf.paragraphs[0]
        p.text = f"Winio  |  Opportunité de revenu complémentaire  |  {current_slide}/{total_slides}"
        p.font.name = "Poppins"
        p.font.size = Pt(10)
        p.font.color.rgb = c_grey

    blank_layout = prs.slide_layouts[6] # completely blank layout

    # ==========================================
    # SLIDE 1: Bienvenue (Dark)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_background(s1, c_charcoal)
    
    # Add Winio/Winelio Logo
    if os.path.exists(logo_on_dark):
        s1.shapes.add_picture(logo_on_dark, Inches(4.666), Inches(1.0), width=Inches(4.0))
    
    # Title
    tf_title = create_text_box(s1, Inches(1.0), Inches(2.3), Inches(11.333), Inches(0.8))
    p_title = tf_title.paragraphs[0]
    p_title.text = "Une opportunité pour développer un revenu complémentaire"
    p_title.font.name = "Poppins"
    p_title.font.size = Pt(32)
    p_title.font.bold = True
    p_title.font.color.rgb = c_white
    p_title.alignment = PP_ALIGN.CENTER
    
    # Text intro
    tf_intro = create_text_box(s1, Inches(1.5), Inches(3.2), Inches(10.333), Inches(1.5))
    p_intro = tf_intro.paragraphs[0]
    p_intro.text = "Merci d'être présents.\nAujourd'hui, nous allons vous présenter une opportunité simple, accessible et compatible avec votre activité actuelle.\nL'objectif n'est pas de remplacer votre emploi, mais de découvrir comment générer un complément de revenu grâce au réseau et aux recommandations."
    p_intro.font.name = "Poppins"
    p_intro.font.size = Pt(16)
    p_intro.font.color.rgb = RGBColor(220, 225, 230)
    p_intro.alignment = PP_ALIGN.CENTER
    p_intro.space_before = Pt(10)
    
    # Hook Box (Rounded rectangle with orange border)
    hook_bg = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.0), Inches(10.333), Inches(1.4))
    hook_bg.fill.solid()
    hook_bg.fill.fore_color.rgb = RGBColor(53, 59, 61) # slightly lighter charcoal
    hook_bg.line.color.rgb = c_orange
    hook_bg.line.width = Pt(2.0)
    
    tf_hook = create_text_box(s1, Inches(1.8), Inches(5.15), Inches(9.733), Inches(1.1))
    p_hook_lbl = tf_hook.paragraphs[0]
    p_hook_lbl.text = "QUESTION D'ACCROCHE :"
    p_hook_lbl.font.name = "Poppins"
    p_hook_lbl.font.size = Pt(11)
    p_hook_lbl.font.bold = True
    p_hook_lbl.font.color.rgb = c_orange
    p_hook_lbl.alignment = PP_ALIGN.CENTER
    p_hook_lbl.space_after = Pt(4)
    
    p_hook = tf_hook.add_paragraph()
    p_hook.text = "Si vous pouviez gagner 200 €, 300 €, 500 € ou plus chaque mois sans changer votre activité principale, qu'est-ce que cela changerait dans votre quotidien ?"
    p_hook.font.name = "Poppins"
    p_hook.font.size = Pt(16)
    p_hook.font.italic = True
    p_hook.font.color.rgb = c_white
    p_hook.alignment = PP_ALIGN.CENTER

    # Bottom orange/amber accent bar
    accent_bar1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.3), Inches(6.666), Inches(0.2))
    accent_bar1.fill.solid()
    accent_bar1.fill.fore_color.rgb = c_orange
    accent_bar1.line.fill.background()
    
    accent_bar2 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.666), Inches(7.3), Inches(6.667), Inches(0.2))
    accent_bar2.fill.solid()
    accent_bar2.fill.fore_color.rgb = c_amber
    accent_bar2.line.fill.background()

    # ==========================================
    # SLIDE 2: Le contexte économique (Light)
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_background(s2, c_light)
    add_slide_header(s2, "Pourquoi chercher un complément de revenu ?")
    
    # Left Column: The 4 drivers
    card_w = Inches(5.2)
    card_h = Inches(0.85)
    gap_y = Inches(0.18)
    start_y = Inches(1.8)
    
    drivers = [
      ("Hausse du coût de la vie", "Inflation constante touchant les dépenses quotidiennes."),
      ("Pouvoir d'achat sous pression", "Des salaires stables face à des coûts de vie croissants."),
      ("Besoin de sécuriser ses finances", "Disposer d'une épargne de précaution ou d'une sécurité."),
      ("Recherche de liberté financière", "Pouvoir financer des envies et projets sans restriction.")
    ]
    
    for idx, (title, desc) in enumerate(drivers):
        curr_y = start_y + idx * (card_h + gap_y)
        # Background box
        box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), curr_y, card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = c_white
        box.line.color.rgb = RGBColor(230, 235, 238)
        box.line.width = Pt(1.0)
        
        # Left orange bar on each box
        v_bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), curr_y + Inches(0.1), Inches(0.06), card_h - Inches(0.2))
        v_bar.fill.solid()
        v_bar.fill.fore_color.rgb = c_orange if idx % 2 == 0 else c_amber
        v_bar.line.fill.background()
        
        # Text
        tf_card = create_text_box(s2, Inches(1.1), curr_y + Inches(0.1), card_w - Inches(0.4), card_h - Inches(0.2))
        pt1 = tf_card.paragraphs[0]
        pt1.text = title
        pt1.font.name = "Poppins"
        pt1.font.size = Pt(14)
        pt1.font.bold = True
        pt1.font.color.rgb = c_charcoal
        
        pt2 = tf_card.add_paragraph()
        pt2.text = desc
        pt2.font.name = "Poppins"
        pt2.font.size = Pt(11)
        pt2.font.color.rgb = c_grey

    # Right Column: Big Message Card
    msg_box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), start_y, Inches(5.7), Inches(4.0))
    msg_box.fill.solid()
    msg_box.fill.fore_color.rgb = c_white
    msg_box.line.color.rgb = c_amber
    msg_box.line.width = Pt(1.5)
    
    tf_msg = create_text_box(s2, Inches(7.2), start_y + Inches(0.4), Inches(4.9), Inches(3.2))
    p_msg_title = tf_msg.paragraphs[0]
    p_msg_title.text = "LE CONTEXTE"
    p_msg_title.font.name = "Poppins"
    p_msg_title.font.size = Pt(12)
    p_msg_title.font.bold = True
    p_msg_title.font.color.rgb = c_amber
    p_msg_title.space_after = Pt(14)
    
    p_msg_body = tf_msg.add_paragraph()
    p_msg_body.text = "Aujourd'hui, de nombreuses personnes développent une double activité pour :\n\n•  Améliorer leur confort de vie\n•  Financer de nouveaux projets\n•  Préparer des vacances de rêve\n•  Gagner en sérénité financière face à l'imprévu."
    p_msg_body.font.name = "Poppins"
    p_msg_body.font.size = Pt(15)
    p_msg_body.font.color.rgb = c_slate
    p_msg_body.space_before = Pt(8)
    
    add_slide_footer(s2, 2)

    # ==========================================
    # SLIDE 3: Chacun ses objectifs (Light)
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_background(s3, c_light)
    add_slide_header(s3, "Quels sont vos objectifs ?")
    
    # 4 objective cards in a 2x2 grid or 1x4 horizontal
    card_w = Inches(2.7)
    card_h = Inches(2.8)
    gap_x = Inches(0.3)
    start_x = Inches(0.8)
    y_pos = Inches(1.8)
    
    objectives = [
        ("200 € / mois", "Alléger son quotidien", "Financer les loisirs, restos et sorties sans toucher au budget principal."),
        ("500 € / mois", "Sérénité mensuelle", "Prendre en charge des factures d'énergie, assurance ou remboursement de crédit."),
        ("1 000 €+ / mois", "Accélérateur de projets", "Financer de grands voyages, travaux ou l'achat d'un nouveau véhicule."),
        ("Long terme", "Activité indépendante", "Bâtir une véritable activité sur le long terme avec des revenus évolutifs.")
    ]
    
    for idx, (amount, title, details) in enumerate(objectives):
        cx = start_x + idx * (card_w + gap_x)
        
        # Draw card
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y_pos, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = c_white
        # Highlight card 2 and 3 as most common
        if idx in [1, 2]:
            card.line.color.rgb = c_orange
            card.line.width = Pt(1.5)
        else:
            card.line.color.rgb = RGBColor(220, 225, 230)
            card.line.width = Pt(1.0)
            
        tf_card = create_text_box(s3, cx + Inches(0.2), y_pos + Inches(0.2), card_w - Inches(0.4), card_h - Inches(0.4))
        
        p_amt = tf_card.paragraphs[0]
        p_amt.text = amount
        p_amt.font.name = "Poppins"
        p_amt.font.size = Pt(22)
        p_amt.font.bold = True
        p_amt.font.color.rgb = c_orange if idx in [1, 2] else c_charcoal
        p_amt.space_after = Pt(4)
        
        p_title = tf_card.add_paragraph()
        p_title.text = title
        p_title.font.name = "Poppins"
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = c_amber
        p_title.space_after = Pt(10)
        
        p_det = tf_card.add_paragraph()
        p_det.text = details
        p_det.font.name = "Poppins"
        p_det.font.size = Pt(11)
        p_det.font.color.rgb = c_grey
        
    # Key Message at the bottom
    tf_bottom = create_text_box(s3, Inches(0.8), Inches(5.0), Inches(11.733), Inches(1.2))
    p_b_title = tf_bottom.paragraphs[0]
    p_b_title.text = "MESSAGE CLÉ :"
    p_b_title.font.name = "Poppins"
    p_b_title.font.size = Pt(12)
    p_b_title.font.bold = True
    p_b_title.font.color.rgb = c_slate
    p_b_title.space_after = Pt(4)
    
    p_b_body = tf_bottom.add_paragraph()
    p_b_body.text = "Il n'existe pas un seul objectif. Chacun avance à son propre rythme, selon ses ambitions et sa disponibilité.\nL'activité s'adapte à votre vie, et non l'inverse."
    p_b_body.font.name = "Poppins"
    p_b_body.font.size = Pt(15)
    p_b_body.font.italic = True
    p_b_body.font.color.rgb = c_charcoal
    
    add_slide_footer(s3, 3)

    # ==========================================
    # SLIDE 4: La solution : Winio (Light)
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_background(s4, c_light)
    add_slide_header(s4, "La solution : Winio")
    
    # Left: Text Content
    tf_left = create_text_box(s4, Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.5))
    p_s4_intro = tf_left.paragraphs[0]
    p_s4_intro.text = "Qu'est-ce que Winio ?"
    p_s4_intro.font.name = "Poppins"
    p_s4_intro.font.size = Pt(20)
    p_s4_intro.font.bold = True
    p_s4_intro.font.color.rgb = c_orange
    p_s4_intro.space_after = Pt(14)
    
    bullets = [
        "Présentation rapide de l'application : Une interface moderne pour connecter les particuliers aux professionnels.",
        "Concept simple et accessible : Pas besoin de compétences techniques complexes, l'application vous guide.",
        "Développement basé sur les recommandations : Valorisez le bouche-à-oreille et votre réseau existant."
    ]
    for idx, b in enumerate(bullets):
        p = tf_left.add_paragraph()
        p.text = "•  " + b
        p.font.name = "Poppins"
        p.font.size = Pt(15)
        p.font.color.rgb = c_slate
        p.space_after = Pt(14)
        
    # Right: Video Placeholder Card
    v_box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.0))
    v_box.fill.solid()
    v_box.fill.fore_color.rgb = c_charcoal
    v_box.line.color.rgb = c_orange
    v_box.line.width = Pt(2.0)
    
    # Video details inside box
    tf_v = create_text_box(s4, Inches(7.7), Inches(2.8), Inches(4.6), Inches(2.0))
    pv_icon = tf_v.paragraphs[0]
    pv_icon.text = "🎥"
    pv_icon.font.name = "Arial"
    pv_icon.font.size = Pt(48)
    pv_icon.alignment = PP_ALIGN.CENTER
    pv_icon.space_after = Pt(8)
    
    pv_lbl = tf_v.add_paragraph()
    pv_lbl.text = "Vidéo de présentation Winio"
    pv_lbl.font.name = "Poppins"
    pv_lbl.font.size = Pt(16)
    pv_lbl.font.bold = True
    pv_lbl.font.color.rgb = c_white
    pv_lbl.alignment = PP_ALIGN.CENTER
    
    pv_sub = tf_v.add_paragraph()
    pv_sub.text = "[Cliquez pour lancer la vidéo]"
    pv_sub.font.name = "Poppins"
    pv_sub.font.size = Pt(12)
    pv_sub.font.color.rgb = c_grey
    pv_sub.alignment = PP_ALIGN.CENTER
    
    add_slide_footer(s4, 4)

    # ==========================================
    # SLIDE 5: Comment ça fonctionne ? (Light)
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_background(s5, c_light)
    add_slide_header(s5, "Un modèle simple et duplicable")
    
    # 4 timeline boxes horizontally
    step_w = Inches(2.7)
    step_h = Inches(2.2)
    gap_x = Inches(0.3)
    start_x = Inches(0.8)
    y_pos = Inches(1.8)
    
    steps = [
        ("1. Utiliser", "Utiliser l'application", "Familiarisez-vous avec l'outil et découvrez la qualité de l'interface."),
        ("2. Partager", "Partager son expérience", "Parlez de l'application et de ses services autour de vous naturellement."),
        ("3. Développer", "Développer son réseau", "Invitez vos connaissances à utiliser Winio pour créer votre communauté."),
        ("4. Accompagner", "Accompagner le réseau", "Guidez les personnes de votre réseau qui souhaitent faire la même chose.")
    ]
    
    for idx, (step_num, title, desc) in enumerate(steps):
        cx = start_x + idx * (step_w + gap_x)
        
        # Arrow shape / box
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y_pos, step_w, step_h)
        box.fill.solid()
        box.fill.fore_color.rgb = c_white
        box.line.color.rgb = c_orange if idx == 2 else RGBColor(220, 225, 230)
        box.line.width = Pt(1.5) if idx == 2 else Pt(1.0)
        
        tf_step = create_text_box(s5, cx + Inches(0.2), y_pos + Inches(0.2), step_w - Inches(0.4), step_h - Inches(0.4))
        p_num = tf_step.paragraphs[0]
        p_num.text = step_num
        p_num.font.name = "Poppins"
        p_num.font.size = Pt(16)
        p_num.font.bold = True
        p_num.font.color.rgb = c_orange
        p_num.space_after = Pt(4)
        
        p_title = tf_step.add_paragraph()
        p_title.text = title
        p_title.font.name = "Poppins"
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = c_charcoal
        p_title.space_after = Pt(8)
        
        p_desc = tf_step.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Poppins"
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = c_grey
        
    # Message clé box at the bottom
    msg_y = Inches(4.5)
    m_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), msg_y, Inches(11.733), Inches(1.5))
    m_box.fill.solid()
    m_box.fill.fore_color.rgb = RGBColor(255, 245, 240) # light orange tint
    m_box.line.color.rgb = c_orange
    m_box.line.width = Pt(1.0)
    
    tf_m = create_text_box(s5, Inches(1.2), msg_y + Inches(0.3), Inches(10.933), Inches(0.9))
    p_m_lbl = tf_m.paragraphs[0]
    p_m_lbl.text = "MESSAGE CLÉ :"
    p_m_lbl.font.name = "Poppins"
    p_m_lbl.font.size = Pt(12)
    p_m_lbl.font.bold = True
    p_m_lbl.font.color.rgb = c_orange
    p_m_lbl.space_after = Pt(4)
    
    p_m_body = tf_m.add_paragraph()
    p_m_body.text = "Pas besoin d'être un expert, un influenceur ou un commercial aguerri. Le système de Winio repose entièrement sur la simplicité et la duplication d'actions faciles à réaliser."
    p_m_body.font.name = "Poppins"
    p_m_body.font.size = Pt(14.5)
    p_m_body.font.color.rgb = c_slate
    
    add_slide_footer(s5, 5)

    # ==========================================
    # SLIDE 6: Jusqu'où pouvez-vous aller ? (Light)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_background(s6, c_light)
    add_slide_header(s6, "Votre potentiel de développement")
    
    # Left: Explanation
    tf_left = create_text_box(s6, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.5))
    p_s6_t = tf_left.paragraphs[0]
    p_s6_t.text = "Le Simulateur de Gains"
    p_s6_t.font.name = "Poppins"
    p_s6_t.font.size = Pt(20)
    p_s6_t.font.bold = True
    p_s6_t.font.color.rgb = c_orange
    p_s6_t.space_after = Pt(14)
    
    bullets = [
        "Présentation du simulateur de gains : Un outil interactif pour modéliser vos revenus potentiels selon vos actions.",
        "Visualisation des différents scénarios : Observez ce que rapporte le parrainage direct et indirect.",
        "Comprendre l'impact de la régularité : La clé de la réussite réside dans de petites actions régulières plutôt que des efforts ponctuels."
    ]
    for idx, b in enumerate(bullets):
        p = tf_left.add_paragraph()
        p.text = "•  " + b
        p.font.name = "Poppins"
        p.font.size = Pt(15)
        p.font.color.rgb = c_slate
        p.space_after = Pt(14)
        
    # Right: Simulator mockup card
    sim_box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.8), Inches(5.333), Inches(4.0))
    sim_box.fill.solid()
    sim_box.fill.fore_color.rgb = c_white
    sim_box.line.color.rgb = c_amber
    sim_box.line.width = Pt(1.5)
    
    tf_sim = create_text_box(s6, Inches(7.5), Inches(2.1), Inches(4.733), Inches(3.4))
    ps_title = tf_sim.paragraphs[0]
    ps_title.text = "SIMULATEUR DE REVENUS"
    ps_title.font.name = "Poppins"
    ps_title.font.size = Pt(12)
    ps_title.font.bold = True
    ps_title.font.color.rgb = c_amber
    ps_title.space_after = Pt(12)
    
    # Mock visual text for simulator
    ps_line1 = tf_sim.add_paragraph()
    ps_line1.text = "Recommandations personnelles :  5 / mois\nPartenaires parrainés (Niv 1) :      3\nPartenaires indirects (Niv 2-5) :   12"
    ps_line1.font.name = "Consolas"
    ps_line1.font.size = Pt(13)
    ps_line1.font.color.rgb = c_slate
    ps_line1.space_after = Pt(16)
    
    ps_res = tf_sim.add_paragraph()
    ps_res.text = "REVENUS ESTIMÉS : 420 € / mois"
    ps_res.font.name = "Poppins"
    ps_res.font.size = Pt(18)
    ps_res.font.bold = True
    ps_res.font.color.rgb = c_orange
    ps_res.space_after = Pt(16)
    
    # Bottom Question box
    ps_q = tf_sim.add_paragraph()
    ps_q.text = "Et si votre réseau travaillait avec vous pour atteindre vos objectifs ?"
    ps_q.font.name = "Poppins"
    ps_q.font.size = Pt(13)
    ps_q.font.italic = True
    ps_q.font.bold = True
    ps_q.font.color.rgb = c_charcoal
    
    add_slide_footer(s6, 6)

    # ==========================================
    # SLIDE 7: Une activité compatible (Light)
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_background(s7, c_light)
    add_slide_header(s7, "Une opportunité flexible")
    
    # Two Columns/Cards
    card_w = Inches(5.6)
    card_h = Inches(3.2)
    card_y = Inches(1.8)
    
    # Left Card: Salarié
    card_sal = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), card_y, card_w, card_h)
    card_sal.fill.solid()
    card_sal.fill.fore_color.rgb = c_white
    card_sal.line.color.rgb = RGBColor(220, 225, 230)
    card_sal.line.width = Pt(1.0)
    
    tf_sal = create_text_box(s7, Inches(1.1), card_y + Inches(0.3), card_w - Inches(0.6), card_h - Inches(0.6))
    ps_t = tf_sal.paragraphs[0]
    ps_t.text = "COMPATIBLE SALARIAT"
    ps_t.font.name = "Poppins"
    ps_t.font.size = Pt(14)
    ps_t.font.bold = True
    ps_t.font.color.rgb = c_orange
    ps_t.space_after = Pt(10)
    
    ps_body = tf_sal.add_paragraph()
    ps_body.text = "•  Zéro conflit d'intérêt : L'apport d'affaires s'effectue sur votre temps libre.\n•  Sécurité d'emploi : Vous conservez votre salaire fixe à 100%.\n•  Revenu additionnel : Pas de pression de chiffre d'affaires, chaque recommandation validée est un bonus."
    ps_body.font.name = "Poppins"
    ps_body.font.size = Pt(13)
    ps_body.font.color.rgb = c_slate
    ps_body.space_before = Pt(8)
    
    # Right Card: Indépendant
    card_ind = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), card_y, card_w, card_h)
    card_ind.fill.solid()
    card_ind.fill.fore_color.rgb = c_white
    card_ind.line.color.rgb = RGBColor(220, 225, 230)
    card_ind.line.width = Pt(1.0)
    
    tf_ind = create_text_box(s7, Inches(7.2), card_y + Inches(0.3), card_w - Inches(0.6), card_h - Inches(0.6))
    pi_t = tf_ind.paragraphs[0]
    pi_t.text = "COMPATIBLE INDÉPENDANT"
    pi_t.font.name = "Poppins"
    pi_t.font.size = Pt(14)
    pi_t.font.bold = True
    pi_t.font.color.rgb = c_amber
    pi_t.space_after = Pt(10)
    
    pi_body = tf_ind.add_paragraph()
    pi_body.text = "•  Synergie d'affaires : Recommandez des services complémentaires à vos propres clients.\n•  Monétisation du réseau : Valorisez les contacts que vous ne pouvez pas traiter directement.\n•  Diversification des revenus : Créez une nouvelle source de revenus récurrents."
    pi_body.font.name = "Poppins"
    pi_body.font.size = Pt(13)
    pi_body.font.color.rgb = c_slate
    pi_body.space_before = Pt(8)

    # Bottom Message Box
    msg_y = Inches(5.3)
    mb = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), msg_y, Inches(11.733), Inches(1.1))
    mb.fill.solid()
    mb.fill.fore_color.rgb = c_white
    mb.line.color.rgb = c_orange
    mb.line.width = Pt(1.5)
    
    tf_mb = create_text_box(s7, Inches(1.1), msg_y + Inches(0.2), Inches(11.133), Inches(0.7))
    pmb = tf_mb.paragraphs[0]
    pmb.text = "Gestion libre de son temps & Développement progressif :"
    pmb.font.name = "Poppins"
    pmb.font.size = Pt(12)
    pmb.font.bold = True
    pmb.font.color.rgb = c_orange
    pmb.space_after = Pt(2)
    
    pmb_sub = tf_mb.add_paragraph()
    pmb_sub.text = "Message clé : Vous gardez votre activité principale et vous avancez entièrement à votre propre rythme."
    pmb_sub.font.name = "Poppins"
    pmb_sub.font.size = Pt(14)
    pmb_sub.font.color.rgb = c_charcoal
    
    add_slide_footer(s7, 7)

    # ==========================================
    # SLIDE 8: La vraie question (Objections - Light)
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_background(s8, c_light)
    add_slide_header(s8, "Qu'est-ce qui vous empêcherait de démarrer ?")
    
    # Left: Grid of 4 Objections (Fears)
    start_y = Inches(1.8)
    card_w = Inches(5.5)
    card_h = Inches(1.0)
    gap_y = Inches(0.2)
    
    objections = [
        ("Manque de temps ?", "L'activité ne requiert que quelques heures par semaine à planifier à votre guise."),
        ("Manque d'expérience ?", "Aucun prérequis. Winio fournit les outils et des modèles pas-à-pas."),
        ("Manque de confiance ?", "Le parrainage signifie que vous n'êtes jamais seul. Vous êtes guidé par votre sponsor."),
        ("Peur de ne pas savoir faire ?", "La recommandation est naturelle : nous recommandons déjà tous des restos ou des films !")
    ]
    
    for idx, (title, desc) in enumerate(objections):
        curr_y = start_y + idx * (card_h + gap_y)
        
        # Draw box
        box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), curr_y, card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = c_white
        box.line.color.rgb = RGBColor(230, 235, 238)
        box.line.width = Pt(1.0)
        
        # Bullet indicator
        ind = s8.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), curr_y + Inches(0.35), Inches(0.3), Inches(0.3))
        ind.fill.solid()
        ind.fill.fore_color.rgb = c_orange
        ind.line.fill.background()
        
        tf_card = create_text_box(s8, Inches(1.4), curr_y + Inches(0.15), card_w - Inches(0.7), card_h - Inches(0.3))
        pt_t = tf_card.paragraphs[0]
        pt_t.text = title
        pt_t.font.name = "Poppins"
        pt_t.font.size = Pt(13)
        pt_t.font.bold = True
        pt_t.font.color.rgb = c_charcoal
        
        pt_d = tf_card.add_paragraph()
        pt_d.text = desc
        pt_d.font.name = "Poppins"
        pt_d.font.size = Pt(11)
        pt_d.font.color.rgb = c_grey

    # Right: Mindset Answer Card
    ans_box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.6))
    ans_box.fill.solid()
    ans_box.fill.fore_color.rgb = c_charcoal
    ans_box.line.color.rgb = c_orange
    ans_box.line.width = Pt(2.0)
    
    tf_ans = create_text_box(s8, Inches(7.2), Inches(2.3), Inches(4.9), Inches(3.6))
    pa_lbl = tf_ans.paragraphs[0]
    pa_lbl.text = "NOTRE VISION"
    pa_lbl.font.name = "Poppins"
    pa_lbl.font.size = Pt(12)
    pa_lbl.font.bold = True
    pa_lbl.font.color.rgb = c_orange
    pa_lbl.space_after = Pt(24)
    
    pa_body = tf_ans.add_paragraph()
    pa_body.text = "Chaque réussite commence par une décision."
    pa_body.font.name = "Poppins"
    pa_body.font.size = Pt(22)
    pa_body.font.bold = True
    pa_body.font.color.rgb = c_white
    pa_body.space_after = Pt(18)
    
    pa_sub = tf_ans.add_paragraph()
    pa_sub.text = "L'importante décision n'est pas de tout maîtriser aujourd'hui, mais d'être simplement prêt à apprendre et à avancer ensemble."
    pa_sub.font.name = "Poppins"
    pa_sub.font.size = Pt(15)
    pa_sub.font.color.rgb = RGBColor(220, 225, 230)
    
    add_slide_footer(s8, 8)

    # ==========================================
    # SLIDE 9: Conclusion & Passage à l'action (Dark)
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    set_background(s9, c_charcoal)
    
    if os.path.exists(logo_on_dark):
        s9.shapes.add_picture(logo_on_dark, Inches(10.5), Inches(0.5), width=Inches(2.0))
        
    # Title
    tf_s9_t = create_text_box(s9, Inches(0.8), Inches(0.5), Inches(9.0), Inches(0.8))
    p_s9_t = tf_s9_t.paragraphs[0]
    p_s9_t.text = "Votre prochaine étape"
    p_s9_t.font.name = "Poppins"
    p_s9_t.font.size = Pt(28)
    p_s9_t.font.bold = True
    p_s9_t.font.color.rgb = c_white
    
    dec_line9 = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.06))
    dec_line9.fill.solid()
    dec_line9.fill.fore_color.rgb = c_orange
    dec_line9.line.fill.background()
    
    # 4 Action Steps layout horizontally
    act_w = Inches(2.7)
    act_h = Inches(1.8)
    gap_x = Inches(0.3)
    start_x = Inches(0.8)
    y_pos = Inches(1.8)
    
    act_steps = [
        ("1. DÉFINIR", "Définissez votre objectif", "Combien souhaitez-vous gagner chaque mois pour démarrer ?"),
        ("2. DÉCOUVRIR", "Découvrez le potentiel", "Explorez l'application, testez le simulateur et posez vos questions."),
        ("3. APPRENDRE", "Faites-vous accompagner", "Intégrez le réseau avec votre parrain pour être guidé pas-à-pas."),
        ("4. DÉMARRER", "Commencez dès aujourd'hui", "Activez votre compte de parrainage et invitez vos premiers contacts.")
    ]
    
    for idx, (label, title, desc) in enumerate(act_steps):
        cx = start_x + idx * (act_w + gap_x)
        
        # Card background
        card = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y_pos, act_w, act_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(53, 59, 61)
        card.line.color.rgb = c_amber if idx == 3 else RGBColor(80, 85, 90)
        card.line.width = Pt(1.5) if idx == 3 else Pt(1.0)
        
        tf_act = create_text_box(s9, cx + Inches(0.15), y_pos + Inches(0.15), act_w - Inches(0.3), act_h - Inches(0.3))
        p_lbl = tf_act.paragraphs[0]
        p_lbl.text = label
        p_lbl.font.name = "Poppins"
        p_lbl.font.size = Pt(12)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = c_amber if idx == 3 else c_orange
        p_lbl.space_after = Pt(6)
        
        p_t = tf_act.add_paragraph()
        p_t.text = title
        p_t.font.name = "Poppins"
        p_t.font.size = Pt(12.5)
        p_t.font.bold = True
        p_t.font.color.rgb = c_white
        p_t.space_after = Pt(6)
        
        p_d = tf_act.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Poppins"
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = c_grey
        
    # Closing Quote
    quote_y = Inches(4.1)
    q_box = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), quote_y, Inches(11.733), Inches(2.0))
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = RGBColor(53, 59, 61)
    q_box.line.color.rgb = c_orange
    q_box.line.width = Pt(1.0)
    
    tf_q = create_text_box(s9, Inches(1.2), quote_y + Inches(0.3), Inches(10.933), Inches(1.4))
    pq1 = tf_q.paragraphs[0]
    pq1.text = "Votre situation actuelle est le résultat des décisions prises hier.\nVotre avenir dépend des décisions que vous prenez aujourd'hui."
    pq1.font.name = "Poppins"
    pq1.font.size = Pt(18)
    pq1.font.bold = True
    pq1.font.italic = True
    pq1.font.color.rgb = c_white
    pq1.alignment = PP_ALIGN.CENTER
    pq1.space_after = Pt(8)
    
    pq2 = tf_q.add_paragraph()
    pq2.text = "Prenez la bonne décision pour vos finances et commencez à bâtir dès maintenant."
    pq2.font.name = "Poppins"
    pq2.font.size = Pt(14)
    pq2.font.color.rgb = c_orange
    pq2.alignment = PP_ALIGN.CENTER
    
    # Bottom brand and page
    tf_f = create_text_box(s9, Inches(0.8), Inches(6.8), Inches(11.733), Inches(0.3))
    pf = tf_f.paragraphs[0]
    pf.text = "Winio  |  Recommandez. Connectez. Développez.  |  9/9"
    pf.font.name = "Poppins"
    pf.font.size = Pt(10)
    pf.font.color.rgb = c_grey

    # Save presentation
    output_path = "/Users/steph/PROJETS/WINELIO/winelio/outputs/Presentation_Winio_Opportunity.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
